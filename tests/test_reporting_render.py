"""Sprint B4 report engine: the PPTX renderer (Launch composition × Core primitives).

The acceptance test for the report/deliverable engine (R2 pilot blocker №1): build
the 8-section Launch Forecast deck from the REAL forecast fixture, wiring Launch's
context to the Core `aurora_reporting` primitives, and assert the produced .pptx is
structurally sound, fonts are embedded, and NO forbidden phrase reached the deck.

Launch is the first consumer of the Core reporting SSOT, so this doubles as the
cross-product acceptance test for those primitives (cone / radar / table / badge /
font-embed) on real data.
"""

from __future__ import annotations

import zipfile

import pytest

from pptx import Presentation
from pptx.util import Emu

from aurora_launch.reporting import copy
from aurora_launch.reporting.context import build_report_context
from aurora_launch.reporting.render_pptx import build_launch_forecast_pptx
from aurora_launch.sample_bundles.report_fixture import build_sample_forecast_fixture


@pytest.fixture(scope="module")
def deck(tmp_path_factory) -> dict:
    ctx = build_report_context(build_sample_forecast_fixture())
    out = tmp_path_factory.mktemp("report") / "launch_forecast.pptx"
    manifest = build_launch_forecast_pptx(
        ctx,
        str(out),
        aurora_version="v0.2.5",
        project_id="test-0001-acceptance",
        date_generated="2026-06-15",
        hash_signature="0" * 64,
    )
    return {"path": str(out), "manifest": manifest, "prs": Presentation(str(out))}


class TestDeckStructure:
    def test_slide_count_covers_eight_sections(self, deck: dict) -> None:
        # cover + exec-headline + key-metrics + proxy + radar + caveats + uncertainty
        # + (cone + weekly) × 3 horizons + methodology + model-card = 15.
        assert len(list(deck["prs"].slides)) == 15

    def test_widescreen_16_9(self, deck: dict) -> None:
        prs = deck["prs"]
        assert round(Emu(prs.slide_width).inches, 2) == 13.33
        assert round(Emu(prs.slide_height).inches, 1) == 7.5

    def test_four_embedded_png_charts(self, deck: dict) -> None:
        # radar + 3 forecast cones, each a Core charts_png primitive embedded as PNG.
        pics = sum(1 for s in deck["prs"].slides for sh in s.shapes if sh.shape_type == 13)
        assert pics == 4

    def test_five_styled_tables(self, deck: dict) -> None:
        # key-metrics + uncertainty + weekly-breakdown × 3.
        tables = sum(1 for s in deck["prs"].slides for sh in s.shapes if sh.has_table)
        assert tables == 5


class TestFontEmbedding:
    def test_inter_lora_embedded_in_ooxml(self, deck: dict) -> None:
        with zipfile.ZipFile(deck["path"]) as z:
            font_parts = [n for n in z.namelist() if n.startswith("ppt/fonts/")]
        # Inter Regular/Bold + Lora Regular/Bold.
        assert len(font_parts) == 4

    def test_manifest_reports_embedded(self, deck: dict) -> None:
        assert deck["manifest"]["fonts_embedded"] == ["Inter", "Lora"]


class TestClientSurfaceHygiene:
    def test_no_forbidden_phrase_in_rendered_deck(self, deck: dict) -> None:
        """The §4.3 forbidden-phrase gate must hold in the ACTUAL artifact, not only
        at context-build time (prove the effect on the layer below the surface)."""
        texts: list[str] = []
        for s in deck["prs"].slides:
            for sh in s.shapes:
                if sh.has_text_frame:
                    texts.append(sh.text_frame.text)
                if sh.has_table:
                    for row in sh.table.rows:
                        for cell in row.cells:
                            texts.append(cell.text)
        assert copy.find_forbidden_phrases("\n".join(texts)) == []


class TestPendingDataLogged:
    def test_channel_and_sensitivity_skipped_not_silently_dropped(self, deck: dict) -> None:
        skipped = deck["manifest"]["skipped"]
        # per-channel forecast path (engine follow-up) → these sections have no data
        # yet and must be reported as skipped, not silently omitted.
        for key in ("forecast_12w", "forecast_26w", "forecast_52w"):
            assert f"{key}.channel_decomposition" in skipped
            assert f"{key}.sensitivity" in skipped


class TestBandDerivation:
    def test_derive_bands_nests_95_80_50(self) -> None:
        from aurora_launch.reporting.render_pptx import _derive_bands

        cone = [{"x": 1, "mean": 100.0, "lo": 50.0, "hi": 150.0}]
        bands = _derive_bands(cone)
        assert set(bands) == {95, 80, 50}
        # 95% is the engine's actual interval; 80/50 nest strictly inside it.
        assert bands[95] == ([50.0], [150.0])
        lo95, hi95 = bands[95][0][0], bands[95][1][0]
        lo80, hi80 = bands[80][0][0], bands[80][1][0]
        lo50, hi50 = bands[50][0][0], bands[50][1][0]
        assert lo95 < lo80 < lo50 < 100.0 < hi50 < hi80 < hi95
