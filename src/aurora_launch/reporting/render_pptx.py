"""Launch Forecast → PPTX renderer (the composition layer for Sprint B4).

This is the Launch-side TRANSLATION layer: it takes the neutral 8-section report
context (:func:`aurora_launch.reporting.context.build_report_context`) and wires it
to the Core ``aurora_reporting`` render primitives. Per the CPI boundary agreed in
the relay loop:

  - **Core** owns the generic primitives (theme, matplotlib→PNG charts, styled
    table, vector tier badge, font embedding) — universal across products.
  - **Launch** owns this composition: which slides, in what order, the RU copy,
    the 8-section launch_forecast template, and methodology-specific choices (e.g.
    deriving 80/50 CI bands from the engine's single 95% interval — a Launch
    methodological assumption, NOT a Core concern).

Primitives are imported from Core submodules today; they will move to the single
``aurora_reporting.primitives`` facade Core is adding — a trivial import swap.

Pending real data (engine follow-up, NOT a renderer gap): channel decomposition
(§5.3 stacked-area) and sensitivity (§5.4 tornado) need a per-channel forecast
path; the orchestrator returns total per-period points. Those slides are skipped
and logged. The uncertainty decomposition (§4.2) is rendered as a styled table
until Core ships the ``pie`` primitive (Batch 2).
"""

from __future__ import annotations

import io
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --- Core aurora_reporting primitives, via the monolith-free `primitives` facade
# (one stable import path; no WeasyPrint/jinja pulled — verified in a clean subprocess).
from aurora_reporting.primitives import (
    AURORA_HYBRID,
    TierVerdict,
    forecast_cone,
    similarity_radar,
    styled_table,
    tier_badge,
    tier_for,
)

from aurora_launch.reporting import copy

# 16:9 widescreen deck.
_SLIDE_W_IN = 13.333
_SLIDE_H_IN = 7.5
_MARGIN_IN = 0.7
_THEME = AURORA_HYBRID


def _rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color.lstrip("#").upper())


def _blank_slide(prs: Presentation):
    """A blank slide with the Aurora cream background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(_THEME.bg_white)
    return slide


def _text(
    slide,
    text: str,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int,
    bold: bool = False,
    color: str | None = None,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    font: str | None = None,
    italic: bool = False,
):
    """Add a single-paragraph text box. Returns the text frame for further runs."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font or _THEME.body
    run.font.color.rgb = _rgb(color or _THEME.text_primary)
    return tf


def _bullets(slide, items: list[str], *, x: float, y: float, w: float, h: float, size: int = 13):
    """A bulleted list (each item one paragraph)."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.size = Pt(size)
        run.font.name = _THEME.body
        run.font.color.rgb = _rgb(_THEME.text_primary)
    return tf


def _accent_rule(slide, *, x: float, y: float, w: float = 1.4, lime: bool = False):
    """A 2pt rule under a title — gold by default; Sacred Lime is signature-only
    (the ONE sanctioned lime use, per the theme's lime-guard)."""
    from pptx.enum.shapes import MSO_SHAPE

    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Pt(2.5))
    rule.fill.solid()
    rule.fill.fore_color.rgb = _rgb(_THEME.sig_lime if lime else _THEME.gold)
    rule.line.fill.background()
    rule.shadow.inherit = False
    return rule


def _section_title(slide, kicker: str, title: str, *, lime: bool = False):
    """Standard section header: small gold kicker + large display title + rule."""
    _text(slide, kicker, x=_MARGIN_IN, y=0.5, w=10.0, h=0.4, size=12, bold=True,
          color=_THEME.gold, font=_THEME.body)
    _text(slide, title, x=_MARGIN_IN, y=0.85, w=11.9, h=0.9, size=28, bold=True,
          color=_THEME.deep[0], font=_THEME.display)
    _accent_rule(slide, x=_MARGIN_IN, y=1.7, lime=lime)


def _footer(slide, ctx_cover: dict[str, Any]):
    """Slide footer: brand · version · project id."""
    parts = [
        "Aurora AI · Launch Forecast",
        ctx_cover.get("aurora_version") or "",
        (ctx_cover.get("project_id") or "")[:8],
    ]
    _text(slide, "   ·   ".join(p for p in parts if p), x=_MARGIN_IN, y=_SLIDE_H_IN - 0.45,
          w=12.0, h=0.3, size=8, color=_THEME.text_secondary)


def _add_png(slide, png: bytes, *, x: float, y: float, w: float | None = None,
             h: float | None = None):
    """Embed PNG bytes (from a Core charts_png primitive) as a picture."""
    kwargs: dict[str, Any] = {}
    if w is not None:
        kwargs["width"] = Inches(w)
    if h is not None:
        kwargs["height"] = Inches(h)
    return slide.shapes.add_picture(io.BytesIO(png), Inches(x), Inches(y), **kwargs)


# ── Section builders ─────────────────────────────────────────────────────────


def _cover(prs, ctx):
    slide = _blank_slide(prs)
    cover = ctx["cover"]
    # Deep background band (premium, lime-guard respected — no lime fill).
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(_THEME.deep[0])
    _text(slide, "AURORA AI", x=_MARGIN_IN, y=0.6, w=6.0, h=0.5, size=16, bold=True,
          color=_THEME.gold, font=_THEME.body)
    _text(slide, cover["recipient_brand"], x=_MARGIN_IN, y=2.6, w=11.9, h=1.2, size=44,
          bold=True, color=_THEME.bg_white, font=_THEME.display)
    _accent_rule(slide, x=_MARGIN_IN, y=3.85, w=2.2, lime=True)  # signature lime
    _text(slide, cover["subtitle"], x=_MARGIN_IN, y=4.05, w=11.0, h=0.5, size=18,
          color=_THEME.gold, font=_THEME.body)
    _text(slide, cover["tagline"], x=_MARGIN_IN, y=4.7, w=10.5, h=1.0, size=13,
          color="#C9D4E0", font=_THEME.body)
    footer = "   ·   ".join(p for p in [
        cover.get("aurora_version") or "",
        cover.get("date_generated") or "",
        f"ID {(cover.get('project_id') or '')[:8]}" if cover.get("project_id") else "",
        f"SHA {(cover.get('hash_signature') or '')[:8]}" if cover.get("hash_signature") else "",
    ] if p)
    _text(slide, footer, x=_MARGIN_IN, y=_SLIDE_H_IN - 0.6, w=12.0, h=0.35, size=9,
          color="#7E92A8", font=_THEME.body)
    return slide


def _executive_headline(prs, ctx):
    slide = _blank_slide(prs)
    es = ctx["executive_summary"]
    _section_title(slide, "РЕЗЮМЕ ДЛЯ РУКОВОДСТВА", "Ключевой прогноз", lime=True)
    # Headline number.
    _text(slide, es["headline"], x=_MARGIN_IN, y=2.4, w=9.3, h=1.6, size=22, bold=True,
          color=_THEME.deep[0], font=_THEME.display)
    _text(slide, es["similarity_one_liner"], x=_MARGIN_IN, y=4.2, w=9.3, h=0.6, size=13,
          color=_THEME.text_secondary)
    _text(slide, es["tier"]["verdict"], x=_MARGIN_IN, y=4.9, w=9.3, h=0.8, size=13,
          color=_THEME.text_primary)
    # Tier badge — RU label (app-side) + Core band colour (Core owns thresholds/colour).
    agg = ctx["proxy_quality"]["radar"]["aggregate"]
    core_tier = tier_for(agg)
    ru_tier = TierVerdict(key=es["tier"]["key"], label=es["tier"]["label"], color=core_tier.color)
    tier_badge(slide, 10.4, 2.4, tier=ru_tier, width_in=2.2, height_in=0.5)
    _footer(slide, ctx["cover"])
    return slide


def _key_metrics(prs, ctx):
    slide = _blank_slide(prs)
    _section_title(slide, "РЕЗЮМЕ ДЛЯ РУКОВОДСТВА", "Ключевые метрики")
    km = ctx["executive_summary"]["key_metrics"]
    headers = ["Горизонт", "Прогноз продаж", "95% ДИ", "Уверенность"]
    rows = [
        [f"{r['period_weeks']} нед.", r["total_display"], f"±{r['ci_pct']:g}%", r["tier_label"]]
        for r in km
    ]
    styled_table(slide, _MARGIN_IN, 2.3, 11.9, headers=headers, rows=rows,
                 highlight_row=0, row_height_in=0.5, font_size=13)
    _text(slide, "С уверенностью 95% продажи за первый период будут в показанном "
          "диапазоне. Aurora раскладывает неопределённость на 4 источника — см. раздел "
          "«Оговорки трансфера».", x=_MARGIN_IN, y=5.4, w=11.9, h=1.0, size=12,
          color=_THEME.text_secondary, italic=True)
    _footer(slide, ctx["cover"])
    return slide


def _proxy_selected(prs, ctx):
    slide = _blank_slide(prs)
    pq = ctx["proxy_quality"]
    _section_title(slide, "КАЧЕСТВО ПРОКСИ", "Выбранный прокси-бренд")
    rows_text = [
        f"Прокси-бренд: {pq['proxy_brand']}",
        f"Категория: {pq.get('proxy_category') or '—'}",
        f"Период данных: {pq.get('proxy_data_period') or '—'}",
        f"Итоговая близость: S = {pq['radar']['aggregate']:g} ({pq['radar']['verdict']})",
    ]
    _bullets(slide, rows_text, x=_MARGIN_IN, y=2.4, w=11.0, h=3.0, size=15)
    _footer(slide, ctx["cover"])
    return slide


def _similarity_radar_slide(prs, ctx):
    slide = _blank_slide(prs)
    radar = ctx["proxy_quality"]["radar"]
    _section_title(slide, "КАЧЕСТВО ПРОКСИ", "Карта близости (6 измерений)")
    dims = radar["dimensions"]  # dict {dim: score} — pivot to Core's parallel seqs
    png = similarity_radar(
        dimensions=list(dims),
        scores=list(dims.values()),
        aggregate=radar["aggregate"],
        title=f"S = {radar['aggregate']:g} ({radar['verdict']})",
        theme=_THEME,
    )
    _add_png(slide, png, x=4.6, y=1.95, h=5.1)
    # Per-dimension scores beside the radar.
    dim_lines = [f"{name}: {score:g}" for name, score in dims.items()]
    _bullets(slide, dim_lines, x=_MARGIN_IN, y=2.4, w=3.6, h=4.0, size=13)
    _footer(slide, ctx["cover"])
    return slide


def _transfer_caveats(prs, ctx):
    slide = _blank_slide(prs)
    tc = ctx["transfer_caveats"]
    _section_title(slide, "ОГОВОРКИ ТРАНСФЕРА", "Что переносится, что — нет")
    col_w = 3.7
    _text(slide, "Переносится (форма)", x=_MARGIN_IN, y=2.1, w=col_w, h=0.4, size=14,
          bold=True, color=_THEME.go)
    _bullets(slide, tc["transfers"], x=_MARGIN_IN, y=2.55, w=col_w, h=3.0, size=12)
    _text(slide, "НЕ переносится", x=_MARGIN_IN + col_w + 0.2, y=2.1, w=col_w, h=0.4,
          size=14, bold=True, color=_THEME.stop)
    _bullets(slide, tc["not_transfers"], x=_MARGIN_IN + col_w + 0.2, y=2.55, w=col_w,
             h=3.0, size=12)
    _text(slide, "Восстанавливается из anchors", x=_MARGIN_IN + 2 * (col_w + 0.2), y=2.1,
          w=col_w, h=0.4, size=14, bold=True, color=_THEME.gold)
    _bullets(slide, tc["reconstructed"], x=_MARGIN_IN + 2 * (col_w + 0.2), y=2.55,
             w=col_w, h=3.0, size=12)
    _text(slide, tc["caveat_text"], x=_MARGIN_IN, y=5.9, w=11.9, h=1.0, size=11,
          color=_THEME.text_secondary, italic=True)
    _footer(slide, ctx["cover"])
    return slide


def _uncertainty(prs, ctx):
    """§4.2 — rendered as a styled table until Core ships the `pie` primitive (Batch 2)."""
    slide = _blank_slide(prs)
    tc = ctx["transfer_caveats"]
    _section_title(slide, "ОГОВОРКИ ТРАНСФЕРА", "Декомпозиция неопределённости")
    unc = tc["uncertainty"]
    label_ru = {
        "proxy": "Прокси (proxy)",
        "transfer": "Трансфер (transfer)",
        "anchor": "Anchors",
        "sampling": "Сэмплинг (sampling)",
    }
    rows = [[label_ru.get(k, k), f"{v * 100:g}%"] for k, v in unc.items()]
    styled_table(slide, _MARGIN_IN, 2.3, 6.0, headers=["Источник", "Доля вариации"],
                 rows=rows, row_height_in=0.5, font_size=13)
    inflation = tc.get("inflation_factor")
    if inflation:
        _text(slide, f"Inflation factor: ×{inflation:g} (verdict Medium → расширение 95% CI)",
              x=_MARGIN_IN, y=5.2, w=11.0, h=0.6, size=12, color=_THEME.text_secondary)
    _footer(slide, ctx["cover"])
    return slide


def _derive_bands(cone: list[dict[str, Any]]) -> dict[int, tuple[list[float], list[float]]]:
    """Launch methodological assumption: derive 80%/50% CI bands from the engine's
    single 95% interval assuming normality (z₉₅=1.96, z₈₀=1.28, z₅₀=0.674), so the
    Core ``forecast_cone`` renders a nested fan. The 95% band is the engine's actual
    output; 80/50 are presentational nestings, not new model outputs.
    """
    z95, z80, z50 = 1.96, 1.2816, 0.6745
    mean = [pt["mean"] for pt in cone]
    half95 = [(pt["hi"] - pt["lo"]) / 2.0 for pt in cone]
    sigma = [h / z95 for h in half95]
    bands: dict[int, tuple[list[float], list[float]]] = {
        95: ([pt["lo"] for pt in cone], [pt["hi"] for pt in cone]),
        80: ([m - z80 * s for m, s in zip(mean, sigma)], [m + z80 * s for m, s in zip(mean, sigma)]),
        50: ([m - z50 * s for m, s in zip(mean, sigma)], [m + z50 * s for m, s in zip(mean, sigma)]),
    }
    return bands


def _forecast_cone_slide(prs, ctx, section_key: str, weeks: int):
    section = ctx[section_key]
    if section is None:
        return None
    slide = _blank_slide(prs)
    _section_title(slide, f"ПРОГНОЗ · {weeks} НЕДЕЛЬ", "Веер прогноза с доверительными интервалами")
    cone = section["cone"]
    periods = [pt["x"] for pt in cone]
    mean = [pt["mean"] for pt in cone]
    png = forecast_cone(
        periods=periods,
        mean=mean,
        bands=_derive_bands(cone),
        ylabel="Продажи, ₽",
        title=None,
        theme=_THEME,
        size_px=(1600, 760),
    )
    _add_png(slide, png, x=_MARGIN_IN, y=1.95, w=11.9)
    _footer(slide, ctx["cover"])
    return slide


def _weekly_breakdown_slide(prs, ctx, section_key: str, weeks: int, *, max_rows: int = 8):
    section = ctx[section_key]
    if section is None:
        return None
    slide = _blank_slide(prs)
    _section_title(slide, f"ПРОГНОЗ · {weeks} НЕДЕЛЬ", "Понедельная разбивка")
    wb = section["weekly_breakdown"][:max_rows]
    headers = ["Неделя", "Среднее, ₽", "Нижняя 95%", "Верхняя 95%"]
    rows = [
        [
            str(r["week"]),
            f"{r['mean']:,.0f}".replace(",", " "),
            f"{r['ci_lower']:,.0f}".replace(",", " "),
            f"{r['ci_upper']:,.0f}".replace(",", " "),
        ]
        for r in wb
    ]
    styled_table(slide, _MARGIN_IN, 2.3, 11.9, headers=headers, rows=rows,
                 row_height_in=0.42, font_size=11)
    if len(section["weekly_breakdown"]) > max_rows:
        _text(slide, f"Показаны первые {max_rows} из {len(section['weekly_breakdown'])} недель "
              "(полная разбивка — в XLSX).", x=_MARGIN_IN, y=6.4, w=11.0, h=0.4, size=10,
              color=_THEME.text_secondary, italic=True)
    _footer(slide, ctx["cover"])
    return slide


def _methodology(prs, ctx):
    slide = _blank_slide(prs)
    method = ctx["methodology"]
    _section_title(slide, "МЕТОДОЛОГИЯ", "Математика и источники")
    formulas = method["formulas"]
    _text(slide, "Adstock:   A_t = X_t + λ · A_(t−1)", x=_MARGIN_IN, y=2.2, w=11.0, h=0.5,
          size=15, font=_THEME.mono_font[0], color=_THEME.deep[0])
    _text(slide, "Hill:   H(x) = β · x^γ / (k^γ + x^γ)", x=_MARGIN_IN, y=2.8, w=11.0, h=0.5,
          size=15, font=_THEME.mono_font[0], color=_THEME.deep[0])
    _text(slide, "Академические источники:", x=_MARGIN_IN, y=3.7, w=11.0, h=0.4, size=13,
          bold=True, color=_THEME.text_primary)
    _bullets(slide, method["references"], x=_MARGIN_IN, y=4.15, w=11.0, h=2.2, size=12)
    _footer(slide, ctx["cover"])
    return slide


def _model_card(prs, ctx):
    slide = _blank_slide(prs)
    method = ctx["methodology"]
    _section_title(slide, "МЕТОДОЛОГИЯ", "Карта модели и воспроизводимость")
    _text(slide, method["posterior_update_reminder"], x=_MARGIN_IN, y=2.2, w=11.9, h=1.2,
          size=13, color=_THEME.text_primary)
    _text(slide, method["cross_reference"], x=_MARGIN_IN, y=3.5, w=11.9, h=1.0, size=12,
          color=_THEME.text_secondary, italic=True)
    cover = ctx["cover"]
    repro = [
        f"Версия Aurora Launch: {cover.get('aurora_version') or '—'}",
        f"Project ID: {cover.get('project_id') or '—'}",
        f"Hash (SHA-256): {cover.get('hash_signature') or '—'}",
        "Methodology Certificate (PDF) прилагается отдельно.",
    ]
    _bullets(slide, repro, x=_MARGIN_IN, y=4.7, w=11.9, h=2.0, size=12)
    _footer(slide, ctx["cover"])
    return slide


def build_launch_forecast_pptx(
    context: dict[str, Any],
    output_path: str,
    *,
    aurora_version: str = "v0.2.5",
    project_id: str | None = None,
    date_generated: str | None = None,
    hash_signature: str | None = None,
    embed_fonts: bool = True,
) -> dict[str, Any]:
    """Render the 8-section Launch Forecast deck to ``output_path``.

    Fills the emit-time cover fields (version / id / date / hash), composes the
    slides from Core primitives, saves the .pptx, then embeds Inter/Lora (OFL) for
    a pixel-stable deck. Returns a small manifest (slide count, skipped sections).
    """
    # Fill emit-time cover fields the context left as None.
    cover = context["cover"]
    cover["aurora_version"] = aurora_version
    cover["project_id"] = project_id
    cover["date_generated"] = date_generated
    cover["hash_signature"] = hash_signature

    prs = Presentation()
    prs.slide_width = Inches(_SLIDE_W_IN)
    prs.slide_height = Inches(_SLIDE_H_IN)

    skipped: list[str] = []
    _cover(prs, context)
    _executive_headline(prs, context)
    _key_metrics(prs, context)
    _proxy_selected(prs, context)
    _similarity_radar_slide(prs, context)
    _transfer_caveats(prs, context)
    _uncertainty(prs, context)
    for key, weeks in (("forecast_12w", 12), ("forecast_26w", 26), ("forecast_52w", 52)):
        if context.get(key) is None:
            skipped.append(key)
            continue
        _forecast_cone_slide(prs, context, key, weeks)
        _weekly_breakdown_slide(prs, context, key, weeks)
        # §5.3 channel decomposition / §5.4 tornado need a per-channel forecast path.
        if context[key].get("channel_decomposition") is None:
            skipped.append(f"{key}.channel_decomposition")
        if context[key].get("sensitivity") is None:
            skipped.append(f"{key}.sensitivity")
    _methodology(prs, context)
    _model_card(prs, context)

    prs.save(output_path)

    embedded: list[str] = []
    if embed_fonts:
        from aurora_reporting.primitives import bundled_font_resolver, embed_brand_fonts

        with __import__("contextlib").suppress(Exception):
            if embed_brand_fonts(output_path, ["Inter", "Lora"], font_resolver=bundled_font_resolver):
                embedded = ["Inter", "Lora"]

    return {
        "output_path": output_path,
        "slide_count": len(prs.slides._sldIdLst),
        "skipped": skipped,
        "fonts_embedded": embedded,
    }


def _main() -> None:
    import json

    from aurora_launch.reporting.context import build_report_context
    from aurora_launch.sample_bundles.report_fixture import build_sample_forecast_fixture

    ctx = build_report_context(build_sample_forecast_fixture())
    manifest = build_launch_forecast_pptx(
        ctx,
        "launch_forecast_sample.pptx",
        aurora_version="v0.2.5",
        project_id="demo-0001-pilot-sample",
        date_generated="2026-06-15",
        hash_signature="0" * 64,
    )
    print(json.dumps(manifest, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    _main()
